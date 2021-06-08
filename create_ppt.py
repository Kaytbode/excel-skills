from pptx import Presentation
from pptx.util import Inches
from db.fetch_all import fetch_data

def slide(query, headings, file_name):
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Excel Skills'

    records = fetch_data(query)

    length = len(headings)

    # slide measurement
    cols = length
    # dynamic generation of rows
    rows = len(records) + 1
    
    left = top = Inches(2.0)
    width = Inches(7.0)
    height = Inches(4.0)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns.width = Inches(1.0)

    # column headings
    for index in range(length):
        table.cell(0, index).text = headings[index]

    row = 1

    # column cells
    for record in records:
        row_length = len(record)
        for index in range(row_length):
            table.cell(row, index).text = str(record[index])
        row += 1

    prs.save(file_name)